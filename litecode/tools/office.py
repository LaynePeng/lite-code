"""办公/生产力工具集：文档、表格、演示、图表、数据分析。

将 lite-code 从代码 Agent 扩展到通用办公场景，让 Agent 能直接产出
docx/xlsx/pptx/pdf 等办公文件，以及进行数据分析和生成图表。

所有依赖包已包含在主依赖中（pyproject.toml dependencies），
`pip install -e .` 时自动安装。
"""
from __future__ import annotations

import csv
import io
import json
import logging
import os
import re
import tempfile
from typing import Any, Dict, List, Optional

from ..core.types import ToolDefinition

logger = logging.getLogger("litecode.tools.office")

# ---------------------------------------------------------------- 依赖检查

_HAS_DOCX: bool = False
_HAS_OPENPYXL: bool = False
_HAS_PPTX: bool = False
_HAS_PANDAS: bool = False
_HAS_MATPLOTLIB: bool = False

try:
    from docx import Document as _DocxDoc
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    _HAS_DOCX = True
except ImportError:
    _DocxDoc = None  # type: ignore

try:
    import openpyxl as _openpyxl

    _HAS_OPENPYXL = True
except ImportError:
    _openpyxl = None  # type: ignore

try:
    from pptx import Presentation as _PptxPresentation
    from pptx.util import Inches as _PptxInches

    _HAS_PPTX = True
except ImportError:
    _PptxPresentation = None  # type: ignore

try:
    import pandas as _pd

    _HAS_PANDAS = True
except ImportError:
    _pd = None  # type: ignore

try:
    import matplotlib
    matplotlib.use("Agg")  # 非交互后端，服务器安全
    import matplotlib.pyplot as _plt

    _HAS_MATPLOTLIB = True
except ImportError:
    _plt = None  # type: ignore


def _missing_dep_msg(pkg: str, tools: str) -> str:
    return (
        f"[Office Tools] 需要安装 {pkg} 才能使用 {tools} 工具。\n"
        f"请运行: pip install lite-code[office]  或   pip install {pkg}"
    )


# ---------------------------------------------------------------- 工具函数


def _ensure_output_dir(workspace: str, subdir: str = ".outputs") -> str:
    """确保输出目录存在，返回绝对路径。"""
    out_dir = os.path.join(os.path.abspath(workspace), subdir)
    os.makedirs(out_dir, exist_ok=True)
    return out_dir


def _safe_filename(name: str) -> str:
    """清理文件名，移除不安全字符。"""
    name = re.sub(r'[<>:"/\\|?*]', "_", name)
    return name.strip() or "output"


# ---------------------------------------------------------------- OfficeTools


class OfficeTools:
    def __init__(self, workspace: Optional[str]) -> None:
        # 桌面应用启动时可能未打开项目（workspace=None），此时回落到用户目录，
        # 真正的产出目录在每次任务时以当前 workspace 为准重建
        self.workspace = os.path.abspath(workspace) if workspace else os.path.expanduser("~")
        self._cjk_font_applied = False

    # ------------------------------------------------------------ 工具定义

    def get_tools(self) -> List[ToolDefinition]:
        return [
            ToolDefinition(
                name="docx_create",
                description=(
                    "根据 Markdown 内容生成 Word (.docx) 文档，支持标题、段落、"
                    "列表、表格、粗体/斜体。返回文件路径。"
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "content": {
                            "type": "string",
                            "description": "Markdown 格式的文档正文",
                        },
                        "filename": {
                            "type": "string",
                            "description": "输出文件名（不含路径，默认 '文档.docx'）",
                        },
                        "title": {
                            "type": "string",
                            "description": "文档标题（文档第一行大标题，可选）",
                        },
                    },
                    "required": ["content"],
                },
            ),
            ToolDefinition(
                name="xlsx_create",
                description=(
                    "根据结构化数据生成 Excel (.xlsx) 表格，支持多 sheet。"
                    "返回文件路径。"
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "data": {
                            "type": "string",
                            "description": "JSON 格式数据。可以是："
                            "1) 对象数组 [{'列名': 值, ...}] 自动生成表头；"
                            "2) 嵌套对象 {\"sheet1\": [...], \"sheet2\": [...]} 多 sheet。",
                        },
                        "filename": {
                            "type": "string",
                            "description": "输出文件名（不含路径，默认 '表格.xlsx'）",
                        },
                    },
                    "required": ["data"],
                },
            ),
            ToolDefinition(
                name="pptx_create",
                description=(
                    "根据结构化内容生成 PowerPoint (.pptx) 演示文稿，支持标题幻灯片"
                    "和正文幻灯片。返回文件路径。"
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "slides": {
                            "type": "string",
                            "description": "JSON 数组，每项为 {title: 幻灯片标题, "
                            "content: Markdown 正文, bullets: [要点列表]（可选）}",
                        },
                        "filename": {
                            "type": "string",
                            "description": "输出文件名（不含路径，默认 '演示文稿.pptx'）",
                        },
                        "title": {
                            "type": "string",
                            "description": "封面标题（可选）",
                        },
                    },
                    "required": ["slides"],
                },
            ),
            ToolDefinition(
                name="pdf_create",
                description=(
                    "将 Markdown 内容生成为 PDF 文件。依赖 reportlab 库。"
                    "返回文件路径。"
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "content": {
                            "type": "string",
                            "description": "Markdown 格式内容",
                        },
                        "filename": {
                            "type": "string",
                            "description": "输出文件名（不含路径，默认 '文档.pdf'）",
                        },
                        "title": {
                            "type": "string",
                            "description": "文档标题（可选）",
                        },
                    },
                    "required": ["content"],
                },
            ),
            ToolDefinition(
                name="data_analyze",
                description=(
                    "对 CSV/JSON 数据进行统计分析，返回分析结果文字。"
                    "支持数据概览、统计描述、分组聚合、排序等。"
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "data": {
                            "type": "string",
                            "description": "CSV 格式文本（第一行为表头）或 JSON 数组字符串",
                        },
                        "instructions": {
                            "type": "string",
                            "description": "分析指令，如 '统计各分组平均值'、'按日期排序'、'描述性统计'",
                        },
                        "output_format": {
                            "type": "string",
                            "enum": ["text", "json", "xlsx"],
                            "description": "输出格式：text（默认，返回文字分析）、json（返回 JSON 串）、xlsx（生成 Excel 文件）",
                        },
                    },
                    "required": ["data", "instructions"],
                },
            ),
            ToolDefinition(
                name="chart_make",
                description=(
                    "根据数据生成图表并保存为图片（PNG），返回图片路径。"
                    "支持柱状图、折线图、饼图、散点图、横向柱状图。"
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "data": {
                            "type": "string",
                            "description": "JSON 格式数据，格式为："
                            "{\"labels\": [\"A\", \"B\", ...], \"values\": [10, 20, ...]}"
                            "或 {\"datasets\": [{\"label\": \"系列1\", \"values\": [...]}, ...]}",
                        },
                        "chart_type": {
                            "type": "string",
                            "enum": ["bar", "line", "pie", "scatter", "barh"],
                            "description": "图表类型：bar(柱状图)、line(折线图)、pie(饼图)、scatter(散点图)、barh(横向柱状图)",
                        },
                        "title": {
                            "type": "string",
                            "description": "图表标题（可选）",
                        },
                        "x_label": {
                            "type": "string",
                            "description": "X 轴标签（可选）",
                        },
                        "y_label": {
                            "type": "string",
                            "description": "Y 轴标签（可选）",
                        },
                        "filename": {
                            "type": "string",
                            "description": "输出文件名（不含路径，默认 'chart.png'）",
                        },
                    },
                    "required": ["data", "chart_type"],
                },
            ),
        ]

    # ------------------------------------------------------------ 执行入口

    async def execute(self, name: str, args: Dict[str, Any]) -> str:
        if name == "docx_create":
            return self._docx_create(args)
        if name == "xlsx_create":
            return self._xlsx_create(args)
        if name == "pptx_create":
            return self._pptx_create(args)
        if name == "pdf_create":
            return self._pdf_create(args)
        if name == "data_analyze":
            return self._data_analyze(args)
        if name == "chart_make":
            return self._chart_make(args)
        raise ValueError(f"Unknown Office Tool: {name}")

    # ------------------------------------------------------------ docx

    def _docx_create(self, args: Dict[str, Any]) -> str:
        if not _HAS_DOCX:
            return _missing_dep_msg("python-docx", "docx_create")

        content = args.get("content", "")
        filename = _safe_filename(args.get("filename", "文档.docx"))
        title = args.get("title", "")

        if not filename.lower().endswith(".docx"):
            filename += ".docx"

        doc = _DocxDoc()

        # 标题
        if title:
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Markdown 转 docx 的简化渲染
        self._md_to_docx(doc, content)

        out_dir = _ensure_output_dir(self.workspace)
        filepath = os.path.join(out_dir, filename)
        doc.save(filepath)

        return f"[Office OK]: 已生成 Word 文档 → {filepath}"

    def _md_to_docx(self, doc, md_text: str) -> None:
        """将 Markdown 文本渲染到 python-docx Document 对象。"""
        from docx.oxml.ns import qn

        lines = md_text.split("\n")
        i = 0
        in_table = False
        table_data: List[List[str]] = []
        table_cols = 0

        while i < len(lines):
            line = lines[i]

            # 表格行（管道语法 | col1 | col2 |）
            if line.strip().startswith("|") and line.strip().endswith("|"):
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                if not in_table:
                    table_data = [cells]
                    table_cols = len(cells)
                    in_table = True
                    # 检查下一行是否为分隔行（|---|---|）
                    if i + 1 < len(lines) and re.match(r"^\|[\s\-:]+\|", lines[i + 1].strip()):
                        i += 1  # 跳过分隔行
                else:
                    if len(cells) <= table_cols:
                        table_data.append(cells)
                i += 1
                continue
            else:
                if in_table and table_data:
                    # 渲染表格
                    if len(table_data) >= 2:
                        table = doc.add_table(rows=len(table_data), cols=table_cols)
                        table.style = "Table Grid"
                        for r_idx, row_data in enumerate(table_data):
                            for c_idx, cell_text in enumerate(row_data):
                                if c_idx < table_cols:
                                    cell = table.rows[r_idx].cells[c_idx]
                                    cell.text = cell_text
                                    if r_idx == 0:
                                        # 表头加粗
                                        for paragraph in cell.paragraphs:
                                            for run in paragraph.runs:
                                                run.bold = True
                    doc.add_paragraph()  # 表后空行
                    table_data = []
                    in_table = False
                    continue

            # 空行
            if not line.strip():
                i += 1
                continue

            # 标题
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading_match:
                level = len(heading_match.group(1))
                text = heading_match.group(2).strip()
                doc.add_heading(text, level=level)
                i += 1
                continue

            # 无序列表
            if re.match(r"^[\s]*[-*+]\s+", line):
                text = re.sub(r"^[\s]*[-*+]\s+", "", line)
                p = doc.add_paragraph(style="List Bullet")
                self._add_styled_run(p, text)
                i += 1
                continue

            # 有序列表
            if re.match(r"^\s*\d+[\.\)]\s+", line):
                text = re.sub(r"^\s*\d+[\.\)]\s+", "", line)
                p = doc.add_paragraph(style="List Number")
                self._add_styled_run(p, text)
                i += 1
                continue

            # 代码块
            if line.strip().startswith("```"):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    code_lines.append(lines[i])
                    i += 1
                code_text = "\n".join(code_lines)
                p = doc.add_paragraph()
                run = p.add_run(code_text)
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                p.paragraph_format.left_indent = Inches(0.3)
                i += 1
                continue

            # 普通段落（支持内联格式）
            p = doc.add_paragraph()
            self._add_styled_run(p, line)
            i += 1

        # 结尾处若还有未渲染的表格
        if in_table and table_data and len(table_data) >= 2:
            table = doc.add_table(rows=len(table_data), cols=table_cols)
            table.style = "Table Grid"
            for r_idx, row_data in enumerate(table_data):
                for c_idx, cell_text in enumerate(row_data):
                    if c_idx < table_cols:
                        table.rows[r_idx].cells[c_idx].text = cell_text

    def _add_styled_run(self, paragraph, text: str) -> None:
        """解析内联 Markdown 格式（粗体、斜体、行内代码）并添加到段落。"""
        # 分割：行内代码 `code`
        parts = re.split(r"(`[^`]+`)", text)
        for part in parts:
            if part.startswith("`") and part.endswith("`"):
                run = paragraph.add_run(part[1:-1])
                run.font.name = "Courier New"
                run.font.size = Pt(9)
            else:
                # 分割粗体 **text** 和斜体 *text*
                sub_parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", part)
                for sp in sub_parts:
                    if sp.startswith("**") and sp.endswith("**"):
                        run = paragraph.add_run(sp[2:-2])
                        run.bold = True
                    elif sp.startswith("*") and sp.endswith("*") and not sp.startswith("**"):
                        run = paragraph.add_run(sp[1:-1])
                        run.italic = True
                    else:
                        paragraph.add_run(sp)

    # ------------------------------------------------------------ xlsx

    def _xlsx_create(self, args: Dict[str, Any]) -> str:
        if not _HAS_OPENPYXL:
            return _missing_dep_msg("openpyxl", "xlsx_create")

        data_str = args.get("data", "")
        filename = _safe_filename(args.get("filename", "表格.xlsx"))

        if not filename.lower().endswith(".xlsx"):
            filename += ".xlsx"

        try:
            data = json.loads(data_str)
        except json.JSONDecodeError as e:
            return f"[Office Error]: data 不是有效的 JSON: {e}"

        wb = _openpyxl.Workbook()
        # 删除默认 sheet
        wb.remove(wb.active)

        if isinstance(data, dict):
            # 多 sheet：{"sheet1": [...], "sheet2": [...]}
            for sheet_name, rows in data.items():
                if not isinstance(rows, list) or not rows:
                    continue
                ws = wb.create_sheet(title=str(sheet_name)[:31])
                self._write_rows_to_sheet(ws, rows)
        elif isinstance(data, list):
            # 单 sheet
            ws = wb.active or wb.create_sheet(title="Sheet1")
            self._write_rows_to_sheet(ws, data)
        else:
            return "[Office Error]: data 必须是 JSON 数组或对象"

        out_dir = _ensure_output_dir(self.workspace)
        filepath = os.path.join(out_dir, filename)
        wb.save(filepath)

        return f"[Office OK]: 已生成 Excel 表格 → {filepath}"

    def _write_rows_to_sheet(self, ws, rows: List[Any]) -> None:
        """将数据行写入 worksheet。第一行对象自动提取表头。"""
        if not rows:
            return

        headers: List[str] = []
        if isinstance(rows[0], dict):
            headers = list(rows[0].keys())
            # 写入表头
            for c, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=c, value=str(h))
                cell.font = _openpyxl.styles.Font(bold=True)
            # 写入数据
            for r, row in enumerate(rows, 2):
                for c, h in enumerate(headers, 1):
                    val = row.get(h, "")
                    if isinstance(val, (list, dict)):
                        val = json.dumps(val, ensure_ascii=False)
                    ws.cell(row=r, column=c, value=val)
        else:
            # 简单列表
            for r, val in enumerate(rows, 1):
                if isinstance(val, (list, tuple)):
                    for c, v in enumerate(val, 1):
                        ws.cell(row=r, column=c, value=v)
                else:
                    ws.cell(row=r, column=1, value=val)

    # ------------------------------------------------------------ pptx

    def _pptx_create(self, args: Dict[str, Any]) -> str:
        if not _HAS_PPTX:
            return _missing_dep_msg("python-pptx", "pptx_create")

        slides_str = args.get("slides", "")
        filename = _safe_filename(args.get("filename", "演示文稿.pptx"))
        title = args.get("title", "")

        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"

        try:
            slides_data = json.loads(slides_str)
        except json.JSONDecodeError as e:
            return f"[Office Error]: slides 不是有效的 JSON: {e}"

        if not isinstance(slides_data, list):
            return "[Office Error]: slides 必须是 JSON 数组"

        prs = _PptxPresentation()

        # 封面
        if title:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = f"生成于 lite-code Office"

        for slide_data in slides_data:
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            bullets = slide_data.get("bullets", [])

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title

            if bullets:
                # 使用占位符中的文本框
                body = slide.placeholders[1]
                text_frame = body.text_frame
                text_frame.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        text_frame.paragraphs[0].text = str(bullet)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = str(bullet)
            elif content:
                body = slide.placeholders[1]
                text_frame = body.text_frame
                text_frame.clear()
                text_frame.paragraphs[0].text = content[:500]

        out_dir = _ensure_output_dir(self.workspace)
        filepath = os.path.join(out_dir, filename)
        prs.save(filepath)

        return f"[Office OK]: 已生成演示文稿 → {filepath}"

    # ------------------------------------------------------------ pdf

    def _pdf_create(self, args: Dict[str, Any]) -> str:
        content = args.get("content", "")
        filename = _safe_filename(args.get("filename", "文档.pdf"))
        title = args.get("title", "")

        if not filename.lower().endswith(".pdf"):
            filename += ".pdf"

        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.platypus import (
                Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
                ListFlowable, ListItem, Preformatted,
            )
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER
        except ImportError:
            return (
                "[Office Tools] 需要安装 reportlab 才能使用 pdf_create 工具。\n"
                "请运行: pip install reportlab"
            )

        out_dir = _ensure_output_dir(self.workspace)
        filepath = os.path.join(out_dir, filename)

        doc = SimpleDocTemplate(filepath, pagesize=A4,
                                topMargin=20*mm, bottomMargin=20*mm,
                                leftMargin=20*mm, rightMargin=20*mm)
        styles = getSampleStyleSheet()
        story: List = []

        if title:
            title_style = ParagraphStyle(
                "Title1", parent=styles["Title"],
                fontSize=24, spaceAfter=12*mm,
                alignment=TA_CENTER,
            )
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 6*mm))

        lines = content.split("\n")
        i = 0
        while i < len(lines):
            line = lines[i]

            if not line.strip():
                story.append(Spacer(1, 3*mm))
                i += 1
                continue

            # 标题
            hm = re.match(r"^(#{1,5})\s+(.+)$", line)
            if hm:
                level = len(hm.group(1))
                text = hm.group(2).strip()
                sz = [22, 18, 15, 13, 11][min(level - 1, 4)]
                h_style = ParagraphStyle(
                    f"Heading{level}", parent=styles["Heading1"],
                    fontSize=sz, spaceBefore=6*mm, spaceAfter=3*mm,
                )
                story.append(Paragraph(text, h_style))
                i += 1
                continue

            # 代码块
            if line.strip().startswith("```"):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    code_lines.append(lines[i])
                    i += 1
                code_text = "\n".join(code_lines)
                code_style = ParagraphStyle(
                    "Code", parent=styles["Code"],
                    fontSize=8, leading=10,
                    leftIndent=6*mm, spaceAfter=3*mm,
                )
                story.append(Preformatted(code_text, code_style))
                i += 1
                continue

            # 无序列表
            if re.match(r"^[\s]*[-*+]\s+", line):
                text = re.sub(r"^[\s]*[-*+]\s+", "", line)
                p = Paragraph(text, styles["Normal"])
                story.append(ListFlowable([ListItem(p)], bulletType="bullet",
                                           leftIndent=30, bulletOffsetY=-2))
                i += 1
                continue

            # 有序列表
            if re.match(r"^\s*\d+[\.\)]\s+", line):
                text = re.sub(r"^\s*\d+[\.\)]\s+", "", line)
                p = Paragraph(text, styles["Normal"])
                story.append(ListFlowable([ListItem(p)], bulletType="1",
                                           leftIndent=30, bulletOffsetY=-2))
                i += 1
                continue

            # 普通段落
            p = Paragraph(line, styles["Normal"])
            story.append(p)
            i += 1

        doc.build(story)
        return f"[Office OK]: 已生成 PDF 文档 → {filepath}"

    # ------------------------------------------------------------ 数据分析

    def _data_analyze(self, args: Dict[str, Any]) -> str:
        if not _HAS_PANDAS:
            return _missing_dep_msg("pandas", "data_analyze")

        data_str = args.get("data", "")
        instructions = args.get("instructions", "").strip()
        output_format = args.get("output_format", "text")

        df: Any = None

        # 尝试解析为 JSON
        try:
            data_json = json.loads(data_str)
            if isinstance(data_json, list):
                df = _pd.DataFrame(data_json)
            elif isinstance(data_json, dict):
                df = _pd.DataFrame([data_json])
        except (json.JSONDecodeError, ValueError):
            pass

        # 尝试解析为 CSV
        if df is None:
            try:
                df = _pd.read_csv(io.StringIO(data_str))
            except Exception:
                pass

        if df is None:
            return "[Office Error]: 无法解析数据。请提供 CSV 格式（第一行表头）或 JSON 数组。"

        result_lines: List[str] = []
        result_lines.append(f"📊 数据分析结果")
        result_lines.append(f"数据形状: {df.shape[0]} 行 × {df.shape[1]} 列")
        result_lines.append(f"列名: {', '.join(str(c) for c in df.columns)}")
        result_lines.append("")

        # 根据指令执行分析
        instr_lower = instructions.lower()

        if any(kw in instr_lower for kw in ["描述", "概览", "统计", "summary", "describe"]):
            result_lines.append("--- 描述性统计 ---")
            desc = df.describe(include="all").to_string()
            result_lines.append(desc)

        if any(kw in instr_lower for kw in ["分组", "group", "agg", "聚合"]):
            result_lines.append("")
            result_lines.append("--- 分组统计 ---")
            # 尝试按第一列分组
            cat_cols = df.select_dtypes(include=["object"]).columns.tolist()
            num_cols = df.select_dtypes(include=["number"]).columns.tolist()
            if cat_cols and num_cols:
                for cc in cat_cols[:2]:
                    for nc in num_cols[:3]:
                        try:
                            grouped = df.groupby(cc)[nc].agg(["mean", "sum", "count"])
                            result_lines.append(f"\n按 {cc} 分组 · {nc} :")
                            result_lines.append(grouped.to_string())
                        except Exception:
                            pass

        if any(kw in instr_lower for kw in ["排序", "sort", "top", "最大的"]):
            result_lines.append("")
            result_lines.append("--- 排序结果 ---")
            num_cols = df.select_dtypes(include=["number"]).columns.tolist()
            if num_cols:
                sorted_df = df.sort_values(by=num_cols[0], ascending=False)
                result_lines.append(sorted_df.head(20).to_string())

        if any(kw in instr_lower for kw in ["空值", "null", "缺失", "missing"]):
            result_lines.append("")
            result_lines.append("--- 缺失值统计 ---")
            null_counts = df.isnull().sum()
            result_lines.append(null_counts.to_string())

        if any(kw in instr_lower for kw in ["相关", "corr", "correlation"]):
            result_lines.append("")
            result_lines.append("--- 相关系数矩阵 ---")
            num_df = df.select_dtypes(include=["number"])
            if num_df.shape[1] >= 2:
                result_lines.append(num_df.corr().to_string())

        if not result_lines[2:]:
            # 默认：显示前几行 + 简单统计
            result_lines.append("--- 前 10 行数据 ---")
            result_lines.append(df.head(10).to_string())
            result_lines.append("")
            result_lines.append("--- 数值列统计 ---")
            num_df = df.select_dtypes(include=["number"])
            if not num_df.empty:
                result_lines.append(num_df.describe().to_string())

        if output_format == "xlsx" and _HAS_OPENPYXL:
            # 输出为 Excel 文件
            filename = "数据分析结果.xlsx"
            out_dir = _ensure_output_dir(self.workspace)
            filepath = os.path.join(out_dir, filename)
            df.to_excel(filepath, index=False, engine="openpyxl")
            result_lines.append(f"\n[Office OK]: 已导出 Excel → {filepath}")
        elif output_format == "json":
            return json.dumps(json.loads(df.head(100).to_json(orient="records")),
                              ensure_ascii=False, indent=2)

        return "\n".join(result_lines)

    # ------------------------------------------------------------ 图表

    def _chart_make(self, args: Dict[str, Any]) -> str:
        if not _HAS_MATPLOTLIB:
            return _missing_dep_msg("matplotlib", "chart_make（需同时安装 pandas）")

        data_str = args.get("data", "")
        chart_type = args.get("chart_type", "bar")
        title = args.get("title", "")
        x_label = args.get("x_label", "")
        y_label = args.get("y_label", "")
        filename = _safe_filename(args.get("filename", "chart.png"))

        if not filename.lower().endswith((".png", ".jpg", ".jpeg", ".svg")):
            filename += ".png"

        try:
            data = json.loads(data_str)
        except json.JSONDecodeError as e:
            return f"[Office Error]: data 不是有效的 JSON: {e}"

        _plt.rcParams["figure.dpi"] = 120
        _plt.rcParams["font.size"] = 11
        _plt.rcParams["axes.unicode_minus"] = False

        # 中文字体适配：按平台选常见中文字体，避免标签显示为方框
        if not self._cjk_font_applied:
            for font in ("Microsoft YaHei", "SimHei", "PingFang SC",
                         "Noto Sans CJK SC", "WenQuanYi Micro Hei", "sans-serif"):
                try:
                    from matplotlib import font_manager
                    matches = font_manager.findfont(
                        font_manager.FontProperties(family=font), fallback_to_default=False
                    )
                    if matches:
                        _plt.rcParams["font.family"] = font
                        break
                except Exception:
                    continue
            self._cjk_font_applied = True

        fig, ax = _plt.subplots(figsize=(10, 6))

        labels = data.get("labels", [])
        datasets = data.get("datasets", [])

        if datasets:
            # 多系列
            for ds in datasets:
                ds_label = ds.get("label", "")
                values = ds.get("values", [])
                if chart_type == "bar":
                    ax.bar(labels, values, label=ds_label, alpha=0.8)
                elif chart_type == "barh":
                    ax.barh(labels, values, label=ds_label, alpha=0.8)
                elif chart_type == "line":
                    ax.plot(labels, values, marker="o", label=ds_label, linewidth=2)
                elif chart_type == "scatter":
                    ax.scatter(labels, values, label=ds_label, s=50, alpha=0.7)
                elif chart_type == "pie":
                    # 饼图只取第一个数据集
                    if ds == datasets[0]:
                        ax.pie(values, labels=labels, autopct="%1.1f%%")
                    break
            if chart_type != "pie":
                ax.legend()
        else:
            # 单系列
            values = data.get("values", [])
            if chart_type == "bar":
                ax.bar(labels, values, color="#4f8cff", alpha=0.8)
            elif chart_type == "barh":
                ax.barh(labels, values, color="#4f8cff", alpha=0.8)
            elif chart_type == "line":
                ax.plot(labels, values, marker="o", color="#4f8cff", linewidth=2)
            elif chart_type == "scatter":
                ax.scatter(labels, values, color="#4f8cff", s=50, alpha=0.7)
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct="%1.1f%%")

        if title:
            ax.set_title(title, fontsize=14, fontweight="bold")
        if x_label and chart_type not in ("pie",):
            ax.set_xlabel(x_label)
        if y_label and chart_type not in ("pie",):
            ax.set_ylabel(y_label)

        if chart_type not in ("pie",):
            _plt.xticks(rotation=30, ha="right")
        _plt.tight_layout()

        out_dir = _ensure_output_dir(self.workspace)
        filepath = os.path.join(out_dir, filename)
        _plt.savefig(filepath, bbox_inches="tight")
        _plt.close(fig)

        return f"[Office OK]: 已生成图表 → {filepath}"